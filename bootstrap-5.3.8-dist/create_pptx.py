from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# Create 16:9 widescreen presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_layout = prs.slide_layouts[6]  # Blank

# Colors
DARK_BG = (15, 23, 42)
MID_BG = (30, 27, 75)
LIGHT_BG = (49, 46, 129)
ACCENT_BLUE = (96, 165, 250)
ACCENT_PURPLE = (192, 132, 252)
ACCENT_INDIGO = (129, 140, 248)
TEXT_PRIMARY = (226, 232, 240)
TEXT_SECONDARY = (148, 163, 184)
TEXT_MUTED = (100, 116, 139)
GOLD = (251, 191, 36)

def add_gradient_bg(slide, c1=DARK_BG, c2=MID_BG, c3=LIGHT_BG):
    background = slide.background
    spPr = background._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    if spPr is None:
        return
    for child in list(spPr):
        spPr.remove(child)
    gradFill = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
    for pos, color in [('0', c1), ('50000', c2), ('100000', c3)]:
        gs = etree.SubElement(gsLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
        gs.set('pos', pos)
        scrgbClr = etree.SubElement(gs, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        scrgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
    lin = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
    lin.set('ang', '2700000')
    lin.set('scaled', '0')

def rgb(color):
    from pptx.dml.color import RGBColor
    return RGBColor(*color)

# ========== SLIDE 1: TITLE ==========
s1 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s1)

# Icon
icon = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.916), Inches(1.8), Inches(1.5), Inches(1.5))
icon.fill.solid()
icon.fill.fore_color.rgb = rgb((99, 102, 241))
icon.fill.fore_color.brightness = -0.8
icon.line.fill.background()

bracket = s1.shapes.add_textbox(Inches(5.916), Inches(2.05), Inches(1.5), Inches(1.0))
tf = bracket.text_frame
p = tf.paragraphs[0]
p.text = "{ }"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(ACCENT_INDIGO)
r.font.name = 'Consolas'

# Title
t = s1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.0))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Infix, Prefix & Postfix"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(48)
r.font.bold = True
r.font.color.rgb = rgb(ACCENT_BLUE)
r.font.name = 'Calibri'

# Subtitle
sub = s1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "Expression conversion for compilers & data structures"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(20)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

# Name & ID
name_box = s1.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4))
tf = name_box.text_frame
p = tf.paragraphs[0]
p.text = "Farhan Shahriar"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

id_box = s1.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.4))
tf = id_box.text_frame
p = tf.paragraphs[0]
p.text = "ID: 253-35-217"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(14)
r.font.color.rgb = rgb(TEXT_MUTED)
r.font.name = 'Calibri'

tags = s1.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4))
tf = tags.text_frame
p = tf.paragraphs[0]
p.text = "Data Structures  •  Algorithms  •  Polish Notation"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(13)
r.font.color.rgb = rgb(TEXT_MUTED)
r.font.name = 'Calibri'

# ========== SLIDE 2: THREE NOTATIONS ==========
s2 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s2)

t = s2.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.333), Inches(0.8))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Three ways to write expressions"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

notations = [
    ("Infix", "A + B", "Operator between operands.\nNeeds parentheses & precedence.", ACCENT_BLUE),
    ("Prefix", "+ A B", "Operator before operands.\nAlso called Polish notation.", ACCENT_PURPLE),
    ("Postfix", "A B +", "Operator after operands.\nAlso called Reverse Polish.", ACCENT_INDIGO),
]
card_w = Inches(3.5)
card_h = Inches(2.8)
start_x = Inches(1.0)
gap = Inches(0.5)

for i, (name, expr, desc, color) in enumerate(notations):
    x = start_x + i * (card_w + gap)
    y = Inches(1.8)
    
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = rgb((255,255,255))
    card.fill.fore_color.brightness = -0.94
    card.line.color.rgb = rgb((255,255,255))
    card.line.color.brightness = -0.9
    
    icon = s2.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.25), y + Inches(0.2), Inches(1.0), Inches(1.0))
    icon.fill.solid()
    icon.fill.fore_color.rgb = rgb(color)
    icon.fill.fore_color.brightness = -0.85
    icon.line.fill.background()
    
    n = s2.shapes.add_textbox(x, y + Inches(1.3), card_w, Inches(0.4))
    tf = n.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Calibri'
    
    e = s2.shapes.add_textbox(x, y + Inches(1.75), card_w, Inches(0.4))
    tf = e.text_frame
    p = tf.paragraphs[0]
    p.text = expr
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.color.rgb = rgb(color)
    r.font.name = 'Consolas'
    
    d = s2.shapes.add_textbox(x + Inches(0.2), y + Inches(2.2), card_w - Inches(0.4), Inches(0.5))
    tf = d.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.color.rgb = rgb(TEXT_SECONDARY)
    r.font.name = 'Calibri'

note = s2.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "All three represent the same mathematical expression — just written differently."
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(15)
r.font.color.rgb = rgb(TEXT_MUTED)
r.font.name = 'Calibri'

# ========== SLIDE 3: WHY CONVERT ==========
s3 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s3)

t = s3.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.333), Inches(0.8))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Why do we convert expressions?"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

bullets = [
    ("No parentheses needed", "Prefix and postfix remove ambiguity without brackets."),
    ("Single-pass evaluation", "Computers scan left-to-right and evaluate instantly."),
    ("Stack-friendly", "Postfix maps directly to push / pop operations."),
    ("Compiler design", "Parsers build syntax trees from postfix order."),
]
y_start = Inches(1.8)
for i, (title, desc) in enumerate(bullets):
    y = y_start + i * Inches(0.9)
    dot = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
    dot.line.fill.background()
    
    t = s3.shapes.add_textbox(Inches(1.8), y, Inches(9.5), Inches(0.4))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = rgb(GOLD)
    r.font.name = 'Calibri'
    
    d = s3.shapes.add_textbox(Inches(1.8), y + Inches(0.35), Inches(9.5), Inches(0.4))
    tf = d.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    r = p.runs[0]
    r.font.size = Pt(15)
    r.font.color.rgb = rgb(TEXT_SECONDARY)
    r.font.name = 'Calibri'

ex_bg = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.5), Inches(8.333), Inches(0.8))
ex_bg.fill.solid()
ex_bg.fill.fore_color.rgb = rgb((15, 23, 42))
ex_bg.line.color.rgb = rgb((99, 102, 241))
ex_bg.line.width = Pt(1)

ex = s3.shapes.add_textbox(Inches(2.7), Inches(5.65), Inches(7.933), Inches(0.5))
tf = ex.text_frame
p = tf.paragraphs[0]
p.text = "Infix: (A + B) * C     →     Postfix: A B + C *"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(18)
r.font.color.rgb = rgb((191, 219, 254))
r.font.name = 'Consolas'

# ========== SLIDE 4: INFIX TO POSTFIX ==========
s4 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s4)

t = s4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.7))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Infix to Postfix conversion"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

note = s4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Expression: (A + B) * C - D"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

steps = [
    ("1", "( A + B ) * C - D", "→ scan left to right"),
    ("2", "A B + * C - D", "→ pop operators by precedence"),
    ("3", "A B + C * - D", "→ * has higher precedence than -"),
    ("4", "A B + C * D -", "→ final postfix expression"),
]
y_start = Inches(1.9)
for i, (num, expr, desc) in enumerate(steps):
    y = y_start + i * Inches(0.85)
    circ = s4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), y + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
    circ.fill.fore_color.brightness = -0.85
    circ.line.fill.background()
    
    num_t = s4.shapes.add_textbox(Inches(1.8), y + Inches(0.08), Inches(0.4), Inches(0.35))
    tf = num_t.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Calibri'
    
    ex = s4.shapes.add_textbox(Inches(2.4), y, Inches(5.0), Inches(0.4))
    tf = ex.text_frame
    p = tf.paragraphs[0]
    p.text = expr
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.color.rgb = rgb(TEXT_PRIMARY if i < 3 else ACCENT_PURPLE)
    r.font.bold = (i == 3)
    r.font.name = 'Consolas'
    
    de = s4.shapes.add_textbox(Inches(2.4), y + Inches(0.35), Inches(7.0), Inches(0.3))
    tf = de.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.color.rgb = rgb(TEXT_MUTED)
    r.font.name = 'Calibri'

tag1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(5.5), Inches(3.5), Inches(0.5))
tag1.fill.solid()
tag1.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
tag1.fill.fore_color.brightness = -0.9
tag1.line.color.rgb = rgb(ACCENT_INDIGO)
tag1.line.color.brightness = -0.7

t1 = s4.shapes.add_textbox(Inches(3.0), Inches(5.55), Inches(3.5), Inches(0.4))
tf = t1.text_frame
p = tf.paragraphs[0]
p.text = "Algorithm: Shunting Yard (Dijkstra)"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(13)
r.font.color.rgb = rgb(ACCENT_INDIGO)
r.font.name = 'Calibri'

tag2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.5), Inches(3.5), Inches(0.5))
tag2.fill.solid()
tag2.fill.fore_color.rgb = rgb(GOLD)
tag2.fill.fore_color.brightness = -0.9
tag2.line.color.rgb = rgb(GOLD)
tag2.line.color.brightness = -0.7

t2 = s4.shapes.add_textbox(Inches(6.8), Inches(5.55), Inches(3.5), Inches(0.4))
tf = t2.text_frame
p = tf.paragraphs[0]
p.text = "Stack: holds operators & parentheses"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(13)
r.font.color.rgb = rgb(GOLD)
r.font.name = 'Calibri'

# ========== SLIDE 5: INFIX TO PREFIX ==========
s5 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s5)

t = s5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.7))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Infix to Prefix conversion"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

note = s5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Expression: (A + B) * C - D"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

steps = [
    ("1", "Reverse the infix string", "D - C * ) B + A ("),
    ("2", "Convert reversed to postfix", "D C B A + * -"),
    ("3", "Reverse the result", "- * + A B C D"),
]
y_start = Inches(2.0)
for i, (num, desc, expr) in enumerate(steps):
    y = y_start + i * Inches(1.0)
    circ = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), y + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
    circ.fill.fore_color.brightness = -0.85
    circ.line.fill.background()
    
    num_t = s5.shapes.add_textbox(Inches(1.8), y + Inches(0.08), Inches(0.4), Inches(0.35))
    tf = num_t.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Calibri'
    
    de = s5.shapes.add_textbox(Inches(2.4), y, Inches(4.5), Inches(0.35))
    tf = de.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    r = p.runs[0]
    r.font.size = Pt(15)
    r.font.color.rgb = rgb(TEXT_SECONDARY)
    r.font.name = 'Calibri'
    
    ex = s5.shapes.add_textbox(Inches(2.4), y + Inches(0.35), Inches(8.0), Inches(0.4))
    tf = ex.text_frame
    p = tf.paragraphs[0]
    p.text = expr
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.color.rgb = rgb(ACCENT_PURPLE if i == 2 else TEXT_PRIMARY)
    r.font.bold = (i == 2)
    r.font.name = 'Consolas'

box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.3), Inches(9.333), Inches(0.8))
box.fill.solid()
box.fill.fore_color.rgb = rgb((255,255,255))
box.fill.fore_color.brightness = -0.96
box.line.color.rgb = rgb((255,255,255))
box.line.color.brightness = -0.92

insight = s5.shapes.add_textbox(Inches(2.2), Inches(5.45), Inches(8.933), Inches(0.5))
tf = insight.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Prefix is read right-to-left during evaluation. The root operator always appears first — making it easy to build expression trees top-down."
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(14)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

# ========== SLIDE 6: POSTFIX EVALUATION ==========
s6 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s6)

t = s6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.7))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Evaluating Postfix"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

note = s6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Expression: 5 3 + 2 *   (equals 16)"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

left_panel = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(5.0), Inches(3.2))
left_panel.fill.solid()
left_panel.fill.fore_color.rgb = rgb((255,255,255))
left_panel.fill.fore_color.brightness = -0.96
left_panel.line.color.rgb = rgb((255,255,255))
left_panel.line.color.brightness = -0.92

left_title = s6.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(5.0), Inches(0.4))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Stack steps"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = rgb(ACCENT_INDIGO)
r.font.name = 'Calibri'

stack_steps = [
    "push 5 → [5]",
    "push 3 → [5, 3]",
    "+ → pop 3, 5 → push 8 → [8]",
    "push 2 → [8, 2]",
    "* → pop 2, 8 → push 16 → [16]",
]
y = Inches(2.5)
for step in stack_steps:
    st = s6.shapes.add_textbox(Inches(1.3), y, Inches(4.4), Inches(0.35))
    tf = st.text_frame
    p = tf.paragraphs[0]
    p.text = step
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Consolas'
    y += Inches(0.45)

right_panel = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.0), Inches(3.2))
right_panel.fill.solid()
right_panel.fill.fore_color.rgb = rgb((255,255,255))
right_panel.fill.fore_color.brightness = -0.96
right_panel.line.color.rgb = rgb((255,255,255))
right_panel.line.color.brightness = -0.92

right_title = s6.shapes.add_textbox(Inches(7.0), Inches(2.05), Inches(5.0), Inches(0.4))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "Expression tree"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = rgb(ACCENT_INDIGO)
r.font.name = 'Calibri'

tree_nodes = [
    (Inches(9.25), Inches(2.6), "*", GOLD, True),
    (Inches(8.25), Inches(3.3), "+", ACCENT_INDIGO, False),
    (Inches(10.25), Inches(3.3), "2", TEXT_PRIMARY, False),
    (Inches(7.75), Inches(4.0), "5", TEXT_PRIMARY, False),
    (Inches(8.75), Inches(4.0), "3", TEXT_PRIMARY, False),
]
for x, y, text, color, is_root in tree_nodes:
    node = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x - Inches(0.25), y - Inches(0.18), Inches(0.5), Inches(0.36))
    node.fill.solid()
    node.fill.fore_color.rgb = rgb(color)
    node.fill.fore_color.brightness = -0.85 if not is_root else -0.8
    node.line.fill.background()
    
    nt = s6.shapes.add_textbox(x - Inches(0.25), y - Inches(0.15), Inches(0.5), Inches(0.3))
    tf = nt.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Consolas'

bot = s6.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.4))
tf = bot.text_frame
p = tf.paragraphs[0]
p.text = "Postfix evaluation is O(n) with a single stack pass — no recursion needed."
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(14)
r.font.color.rgb = rgb(TEXT_MUTED)
r.font.name = 'Calibri'

# ========== SLIDE 7: PREFIX EVALUATION ==========
s7 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s7)

t = s7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.7))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Evaluating Prefix"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

note = s7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Expression: * + 5 3 2   (equals 16)"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.color.rgb = rgb(TEXT_SECONDARY)
r.font.name = 'Calibri'

steps = [
    ("1", "Scan right-to-left", "2  3  5  +  *"),
    ("2", "Push operands onto stack", "[2, 3, 5]"),
    ("3", "On operator: pop two, apply, push result", "+ 5 3 → 8 → [2, 8]"),
    ("4", "Continue until one value remains", "* 8 2 → 16 → [16]"),
]
y_start = Inches(1.9)
for i, (num, desc, expr) in enumerate(steps):
    y = y_start + i * Inches(0.9)
    circ = s7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
    circ.fill.fore_color.brightness = -0.85
    circ.line.fill.background()
    
    num_t = s7.shapes.add_textbox(Inches(1.5), y + Inches(0.08), Inches(0.4), Inches(0.35))
    tf = num_t.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = rgb(TEXT_PRIMARY)
    r.font.name = 'Calibri'
    
    de = s7.shapes.add_textbox(Inches(2.1), y, Inches(4.5), Inches(0.35))
    tf = de.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    r = p.runs[0]
    r.font.size = Pt(15)
    r.font.color.rgb = rgb(TEXT_SECONDARY)
    r.font.name = 'Calibri'
    
    ex = s7.shapes.add_textbox(Inches(2.1), y + Inches(0.35), Inches(8.0), Inches(0.4))
    tf = ex.text_frame
    p = tf.paragraphs[0]
    p.text = expr
    r = p.runs[0]
    r.font.size = Pt(17)
    r.font.color.rgb = rgb(ACCENT_PURPLE if i == 3 else TEXT_PRIMARY)
    r.font.bold = (i == 3)
    r.font.name = 'Consolas'

tag1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(5.5), Inches(3.5), Inches(0.5))
tag1.fill.solid()
tag1.fill.fore_color.rgb = rgb(ACCENT_INDIGO)
tag1.fill.fore_color.brightness = -0.9
tag1.line.color.rgb = rgb(ACCENT_INDIGO)
tag1.line.color.brightness = -0.7

t1 = s7.shapes.add_textbox(Inches(3.0), Inches(5.55), Inches(3.5), Inches(0.4))
tf = t1.text_frame
p = tf.paragraphs[0]
p.text = "Scan direction: Right → Left"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(13)
r.font.color.rgb = rgb(ACCENT_INDIGO)
r.font.name = 'Calibri'

tag2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.5), Inches(3.5), Inches(0.5))
tag2.fill.solid()
tag2.fill.fore_color.rgb = rgb(ACCENT_PURPLE)
tag2.fill.fore_color.brightness = -0.9
tag2.line.color.rgb = rgb(ACCENT_PURPLE)
tag2.line.color.brightness = -0.7

t2 = s7.shapes.add_textbox(Inches(6.8), Inches(5.55), Inches(3.5), Inches(0.4))
tf = t2.text_frame
p = tf.paragraphs[0]
p.text = "First pop = right operand"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(13)
r.font.color.rgb = rgb(ACCENT_PURPLE)
r.font.name = 'Calibri'

# ========== SLIDE 8: SUMMARY ==========
s8 = prs.slides.add_slide(slide_layout)
add_gradient_bg(s8)

t = s8.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.7))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Quick reference"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = rgb(TEXT_PRIMARY)
r.font.name = 'Calibri'

headers = ["Property", "Infix", "Prefix", "Postfix"]
header_colors = [TEXT_SECONDARY, ACCENT_BLUE, ACCENT_PURPLE, ACCENT_INDIGO]
rows = [
    ["Operator position", "Middle", "Before", "After"],
    ["Parentheses", "Required", "Not needed", "Not needed"],
    ["Scan direction", "—", "Right → Left", "Left → Right"],
    ["Evaluation", "Complex", "Stack", "Stack"],
    ["Also called", "Standard", "Polish", "Reverse Polish"],
]

table_left = Inches(1.5)
table_top = Inches(1.3)
col_widths = [Inches(3.0), Inches(2.5), Inches(2.5), Inches(2.5)]
row_height = Inches(0.55)

for col, (header, color) in enumerate(zip(headers, header_colors)):
    x = table_left + sum(col_widths[:col])
    cell = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, col_widths[col], row_height)
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb((255,255,255))
    cell.fill.fore_color.brightness = -0.92
    cell.line.color.rgb = rgb((255,255,255))
    cell.line.color.brightness = -0.88
    
    ht = s8.shapes.add_textbox(x, table_top + Inches(0.1), col_widths[col], Inches(0.35))
    tf = ht.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    r.font.name = 'Calibri'

for row_idx, row_data in enumerate(rows):
    y = table_top + row_height * (row_idx + 1)
    for col, cell_text in enumerate(row_data):
        x = table_left + sum(col_widths[:col])
        cell = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[col], row_height)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb((255,255,255))
        cell.fill.fore_color.brightness = -0.96
        cell.line.color.rgb = rgb((255,255,255))
        cell.line.color.brightness = -0.92
        
        ct = s8.shapes.add_textbox(x, y + Inches(0.1), col_widths[col], Inches(0.35))
        tf = ct.text_frame
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(13)
        r.font.color.rgb = rgb(TEXT_PRIMARY)
        r.font.name = 'Consolas' if col > 0 else 'Calibri'

y_cards = Inches(5.0)
card_data = [
    ("A + B * C", "Infix", ACCENT_BLUE),
    ("+ A * B C", "Prefix", ACCENT_PURPLE),
    ("A B C * +", "Postfix", ACCENT_INDIGO),
]
for i, (expr, label, color) in enumerate(card_data):
    x = Inches(1.5 + i * 3.5)
    card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_cards, Inches(3.0), Inches(1.0))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb((255,255,255))
    card.fill.fore_color.brightness = -0.94
    card.line.color.rgb = rgb(color)
    card.line.color.brightness = -0.7
    
    et = s8.shapes.add_textbox(x, y_cards + Inches(0.15), Inches(3.0), Inches(0.4))
    tf = et.text_frame
    p = tf.paragraphs[0]
    p.text = expr
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.color.rgb = rgb(color)
    r.font.name = 'Consolas'
    
    lt = s8.shapes.add_textbox(x, y_cards + Inches(0.55), Inches(3.0), Inches(0.3))
    tf = lt.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = rgb(TEXT_SECONDARY)
    r.font.name = 'Calibri'

# ========== SAVE ==========
prs.save('Infix_Prefix_Postfix_Farhan_Shahriar_253-35-217.pptx')
print("Saved: Infix_Prefix_Postfix_Farhan_Shahriar_253-35-217.pptx")